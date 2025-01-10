from typing import Optional, List, Dict, Any
from datetime import datetime
import uuid
import os
import pickle
import numpy as np
from pydantic import BaseModel
from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.responses import JSONResponse
import faiss
import logging
import random
import string
from sentence_transformers import SentenceTransformer
import io
from typing import List
import textwrap
from groq import Groq

# Additional imports for file processing
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Configuration
GROQ_API_KEY = "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"  # Replace with your actual Groq API key
MODEL_NAME = "mixtral-8x7b-32768"
FAISS_INDEX_PATH = "company_indexes/"
TOKEN_LENGTH = 6
MAX_FILE_SIZE = 1024 * 1024 * 20  # 20MB limit
CHUNK_SIZE = 500
CHUNK_OVERLAP = 100

# Create FastAPI app
app = FastAPI(title="Company Documents API")

# Initialize Groq client
groq_client = Groq(api_key=GROQ_API_KEY)

# Initialize SentenceTransformer for embeddings
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
dimension = 384  # Dimension for all-MiniLM-L6-v2 embeddings

# Storage structures
company_document_stores = {}
company_faiss_indexes = {}
company_tokens = {}
token_to_company = {}
company_files = {}

# Ensure directories exist
os.makedirs(FAISS_INDEX_PATH, exist_ok=True)

class CompanyCreate(BaseModel):
    name: str

class Company(BaseModel):
    id: str
    name: str
    token: str
    created_at: datetime

class FileMetadata(BaseModel):
    filename: str
    upload_date: datetime
    file_size: int
    chunk_count: int

def generate_company_token() -> str:
    """Generate a unique company token."""
    while True:
        token = ''.join(random.choices(string.ascii_uppercase + string.digits, k=TOKEN_LENGTH))
        if token not in token_to_company:
            return token

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into overlapping chunks."""
    chunks = []
    words = text.split()
    
    if not words:
        return chunks
        
    chunk_words = chunk_size
    overlap_words = overlap
    
    for i in range(0, len(words), chunk_words - overlap_words):
        chunk = ' '.join(words[i:i + chunk_words])
        if chunk:
            chunks.append(chunk.strip())
            
    return chunks

def extract_text_from_pdf(file) -> str:
    """Extract text from PDF file."""
    try:
        reader = PdfReader(file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        raise

def extract_text_from_docx(file) -> str:
    """Extract text from DOCX file."""
    try:
        doc = Document(file)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {str(e)}")
        raise

def extract_text_from_pptx(file) -> str:
    """Extract text from PPTX file."""
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {str(e)}")
        raise

async def process_file_content(file: UploadFile) -> List[str]:
    """Process file and extract chunks based on file type."""
    try:
        # Check file size
        file_contents = await file.read()
        if len(file_contents) > MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail="File too large")

        # Determine file type and extract text
        file_type = file.filename.split('.')[-1].lower()
        file_bytes = io.BytesIO(file_contents)

        if file_type == 'pdf':
            text = extract_text_from_pdf(file_bytes)
        elif file_type == 'docx':
            text = extract_text_from_docx(file_bytes)
        elif file_type == 'pptx':
            text = extract_text_from_pptx(file_bytes)
        else:
            # Assume plain text for other file types
            text = file_contents.decode('utf-8')

        # Return empty list if no text was extracted
        if not text.strip():
            return []

        # Chunk the extracted text
        return chunk_text(text)
    except Exception as e:
        logger.error(f"Error processing file content: {str(e)}")
        raise

def get_embedding(text: str) -> np.ndarray:
    """Generate embedding for text using SentenceTransformer."""
    try:
        embedding = embedding_model.encode(text, convert_to_numpy=True)
        return embedding
    except Exception as e:
        logger.error(f"Error generating embedding: {str(e)}")
        raise

async def get_completion(prompt: str) -> str:
    """Generate AI completion using Groq API."""
    try:
        chat_completion = groq_client.chat.completions.create(
            messages=[
                {"role": "system", "content": "You are a helpful assistant that provides concise and accurate answers based on the given context."},
                {"role": "user", "content": prompt}
            ],
            model=MODEL_NAME,
            max_tokens=500
        )
        return chat_completion.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Error generating completion: {str(e)}")
        raise

def init_company_faiss_index(company_id: str):
    """Initialize or load existing FAISS index for company."""
    try:
        index_path = os.path.join(FAISS_INDEX_PATH, f"{company_id}_index.pkl")
        if os.path.exists(index_path):
            with open(index_path, 'rb') as f:
                company_faiss_indexes[company_id] = pickle.load(f)
                company_document_stores[company_id] = pickle.load(f)
                company_files[company_id] = pickle.load(f)
        else:
            company_faiss_indexes[company_id] = faiss.IndexFlatL2(dimension)
            company_document_stores[company_id] = {}
            company_files[company_id] = {}
            save_company_faiss_index(company_id)
    except Exception as e:
        logger.error(f"Error initializing FAISS index: {str(e)}")
        raise

def save_company_faiss_index(company_id: str):
    """Save company's FAISS index and related data to disk."""
    try:
        index_path = os.path.join(FAISS_INDEX_PATH, f"{company_id}_index.pkl")
        with open(index_path, 'wb') as f:
            pickle.dump(company_faiss_indexes[company_id], f)
            pickle.dump(company_document_stores[company_id], f)
            pickle.dump(company_files[company_id], f)
    except Exception as e:
        logger.error(f"Error saving FAISS index: {str(e)}")
        raise

def verify_token(token: str) -> str:
    """Verify token and return company_id."""
    if token not in token_to_company:
        raise HTTPException(status_code=401, detail="Invalid company token")
    return token_to_company[token]

@app.post("/api/companies/register")
async def register_company(company_data: CompanyCreate):
    """Register a new company and return access token."""
    try:
        company_id = str(uuid.uuid4())
        company_token = generate_company_token()
        
        company = Company(
            id=company_id,
            name=company_data.name,
            token=company_token,
            created_at=datetime.utcnow()
        )
        
        company_tokens[company_id] = company_token
        token_to_company[company_token] = company_id
        init_company_faiss_index(company_id)
        
        return {
            "company_token": company_token,
            "message": f"Company registered successfully. Keep this token safe: {company_token}"
        }
    except Exception as e:
        logger.error(f"Error registering company: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/documents/upload")
async def upload_documents(
    token: str = Form(...),
    files: List[UploadFile] = File(...)
):
    """Upload and process documents for a company."""
    company_id = verify_token(token)
    
    try:
        logger.info(f"Processing {len(files)} files for company: {company_id}")
        
        processed_files = []
        for file in files:
            try:
                chunks = await process_file_content(file)
                
                if not chunks:
                    processed_files.append({
                        "filename": file.filename,
                        "error": "No text content extracted"
                    })
                    continue
                
                # Store file metadata
                file_metadata = FileMetadata(
                    filename=file.filename,
                    upload_date=datetime.utcnow(),
                    file_size=0,
                    chunk_count=len(chunks)
                )
                company_files[company_id][file.filename] = file_metadata
                
                # Process chunks and generate embeddings
                chunk_embeddings = []
                chunk_ids = []
                
                for chunk in chunks:
                    embedding = get_embedding(chunk)
                    chunk_embeddings.append(embedding)
                    
                    chunk_id = str(uuid.uuid4())
                    chunk_ids.append(chunk_id)
                    
                    company_document_stores[company_id][chunk_id] = {
                        "id": chunk_id,
                        "parent_file": file.filename,
                        "content": chunk,
                        "metadata": {
                            "upload_timestamp": datetime.utcnow().isoformat()
                        }
                    }
                
                if chunk_embeddings:
                    company_faiss_indexes[company_id].add(
                        np.array(chunk_embeddings, dtype=np.float32)
                    )
                
                processed_files.append({
                    "filename": file.filename,
                    "chunks_processed": len(chunks),
                    "chunk_ids": chunk_ids
                })
                
                logger.info(f"Successfully processed {file.filename} into {len(chunks)} chunks")
                
            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {str(e)}")
                processed_files.append({
                    "filename": file.filename,
                    "error": str(e)
                })
        
        save_company_faiss_index(company_id)
        
        return {
            "message": "Documents processed successfully",
            "processed_files": processed_files
        }
        
    except Exception as e:
        logger.error(f"Error in upload_documents: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/documents/list")
async def list_documents(token: str):
    """List all documents uploaded for a company."""
    company_id = verify_token(token)
    
    try:
        if company_id not in company_files:
            return {"files": []}
        
        files = []
        for filename, metadata in company_files[company_id].items():
            files.append({
                "filename": filename,
                "upload_date": metadata.upload_date.isoformat(),
                "chunk_count": metadata.chunk_count
            })
        
        return {
            "company_id": company_id,
            "files": files
        }
    except Exception as e:
        logger.error(f"Error listing documents: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/query")
async def query_documents(
    token: str,
    query: str
):
    """Query documents and get relevant answers."""
    company_id = verify_token(token)
    
    try:
        # Check if company has any documents
        if not company_files[company_id]:
            return {"answer": "No documents available to search."}
        
        # Generate query embedding
        query_embedding = get_embedding(query)
        
        # Search for similar chunks
        k = 5  # Number of chunks to retrieve
        distances, indices = company_faiss_indexes[company_id].search(
            np.array([query_embedding], dtype=np.float32), 
            k
        )
        
        # Get relevant chunks
        chunks = []
        for idx in indices[0]:
            if idx != -1:  # Valid index
                chunk_id = list(company_document_stores[company_id].keys())[idx]
                chunk = company_document_stores[company_id][chunk_id]
                chunks.append(chunk)
        
        if chunks:
            # Prepare context from relevant chunks
            context = "\n\n".join([chunk["content"] for chunk in chunks])
            
            # Generate answer using Groq
            prompt = f"""Based on the following context, provide a clear and concise answer to the query: "{query}"

Context:
{context}

Answer:"""
            
            answer = await get_completion(prompt)
        else:
            answer = "No relevant information found in the documents."
        
        return {
            "answer": answer,
            "chunks_found": len(chunks)
        }
        
    except Exception as e:
        logger.error(f"Error processing query: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)
